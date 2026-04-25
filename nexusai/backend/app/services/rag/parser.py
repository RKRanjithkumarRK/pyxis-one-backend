"""File content extraction — PDF, DOCX, XLSX, PPTX, TXT, MD, HTML, CSV."""
from __future__ import annotations
import csv
import io
import logging
from pathlib import Path

logger = logging.getLogger("nexusai.rag.parser")

SUPPORTED_TYPES = {"pdf", "docx", "xlsx", "pptx", "txt", "md", "html", "csv"}


def file_type_from_name(filename: str) -> str:
    return Path(filename).suffix.lstrip(".").lower()


def extract_text(content: bytes, file_type: str) -> str:
    """Return plain text from raw file bytes. Raises ValueError on unsupported type."""
    ft = file_type.lower().lstrip(".")
    if ft == "pdf":
        return _extract_pdf(content)
    elif ft == "docx":
        return _extract_docx(content)
    elif ft == "xlsx":
        return _extract_xlsx(content)
    elif ft == "pptx":
        return _extract_pptx(content)
    elif ft in ("txt", "md"):
        return content.decode("utf-8", errors="replace")
    elif ft == "html":
        return _extract_html(content)
    elif ft == "csv":
        return _extract_csv(content)
    else:
        raise ValueError(f"Unsupported file type: {ft}")


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def _extract_html(content: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "meta", "link", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extract_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        rows.append("\t".join(row))
    return "\n".join(rows)
